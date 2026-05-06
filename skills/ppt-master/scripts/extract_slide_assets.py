#!/usr/bin/env python3
"""
Extract Slide Assets v2 — Tách TOÀN BỘ text + hình ảnh từ PPTX/DOCX/PDF

Cải tiến:
- Trích xuất text từ TẤT CẢ shape types (text frame, table, group, chart)
- Nhận diện title, content, bullet points, table data riêng biệt
- Xử lý grouped shapes recursively
- Detect slide layout chính xác hơn

Usage:
    python3 extract_slide_assets.py <input_file> -o <output_dir>
"""

import sys
import os
import json
import re
import shutil
from pathlib import Path

_SCRIPTS_DIR = Path(__file__).resolve().parent
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

# Windows console encoding fix
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")


def _extract_text_from_shape(shape) -> list:
    """Extract all text from a single shape, handling all types."""
    texts = []

    # 1. Text frame (most common)
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            line = para.text.strip()
            if line:
                # Detect bullet level
                level = para.level if para.level else 0
                prefix = "  " * level
                texts.append(f"{prefix}{line}")

    # 2. Table
    if shape.has_table:
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", " ")
                row_cells.append(cell_text)
            if any(c for c in row_cells):
                texts.append(" | ".join(row_cells))

    # 3. Group shape (recursive)
    if shape.shape_type == 6:  # Group shape
        try:
            for child in shape.shapes:
                texts.extend(_extract_text_from_shape(child))
        except Exception:
            pass

    # 4. Chart
    if shape.has_chart:
        try:
            chart = shape.chart
            chart_title = chart.chart_title.text_frame.text if chart.has_title else ""
            if chart_title:
                texts.append(f"[Chart: {chart_title}]")
            # Try to get chart data
            for series in chart.series:
                for point in series.points:
                    texts.append(f"  {point.value}")
        except Exception:
            texts.append("[Chart: data not extractable]")

    return texts


def _detect_slide_type(title: str, texts: list, slide_idx: int, total_slides: int) -> str:
    """Detect slide type based on content analysis."""
    title_lower = title.lower() if title else ""
    all_text = " ".join(texts).lower()

    if slide_idx == 1:
        return "title_slide"
    if slide_idx == total_slides:
        return "closing"

    # Check for data-heavy slides
    table_indicators = ["|", "tong ket", "tong hop", "sla", "uptime", "rule", "alert"]
    if any(ind in all_text for ind in table_indicators):
        if sum(1 for t in texts if "|" in t) >= 3:
            return "data"

    # Check for image-heavy slides
    if len(texts) <= 2:
        return "image_heavy"

    return "content"


def extract_from_pptx(input_path: Path, output_dir: Path) -> dict:
    """Extract text + images from PPTX file with full content detection."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(str(input_path))
    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    slides_data = []
    all_text = []
    total_slides = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts = []
        slide_images = []
        slide_tables = []

        # Extract from ALL shapes
        for shape in slide.shapes:
            # Get text
            shape_texts = _extract_text_from_shape(shape)
            slide_texts.extend(shape_texts)

            # Get images
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    ext = image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    img_name = f"slide_{slide_idx:02d}_image_{len(slide_images)+1:02d}.{ext}"
                    img_path = images_dir / img_name
                    with open(img_path, "wb") as f:
                        f.write(image.blob)
                    slide_images.append(img_name)
                except Exception:
                    pass

        # Clean and deduplicate texts
        seen = set()
        clean_texts = []
        for t in slide_texts:
            t_clean = t.strip()
            if t_clean and t_clean not in seen:
                seen.add(t_clean)
                clean_texts.append(t_clean)

        # Separate title from content
        title = clean_texts[0] if clean_texts else f"Slide {slide_idx}"
        content = "\n".join(clean_texts)

        # Detect slide type
        slide_type = _detect_slide_type(title, clean_texts, slide_idx, total_slides)

        # Extract bullet points
        bullet_points = []
        for t in clean_texts[1:]:  # Skip title
            if t.startswith("  ") or t.startswith("- ") or t.startswith("* "):
                bullet_points.append(t.strip())
            elif "|" in t:
                slide_tables.append(t)

        slides_data.append({
            "slide_number": slide_idx,
            "title": title,
            "content_raw": content,
            "bullet_points": bullet_points if bullet_points else [],
            "table_data": slide_tables if slide_tables else [],
            "images_extracted": slide_images,
            "layout_hint": slide_type,
            "total_shapes": len(list(slide.shapes)),
            "text_length": len(content),
        })

        # Build markdown
        md_section = f"## Slide {slide_idx}: {title}\n\n{content}"
        if slide_tables:
            md_section += "\n\n### Data:\n" + "\n".join(slide_tables)
        all_text.append(md_section)

    # Write markdown
    md_content = "\n\n---\n\n".join(all_text)
    (output_dir / "extracted.md").write_text(md_content, encoding="utf-8")

    # Write metadata
    (output_dir / "slide_metadata.json").write_text(
        json.dumps(slides_data, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    # Print summary
    total_images = sum(len(s["images_extracted"]) for s in slides_data)
    total_chars = sum(s["text_length"] for s in slides_data)
    print(f"[extract] {len(slides_data)} slides, {total_images} images, {total_chars} chars total")

    return {"slides": slides_data, "total_slides": len(slides_data), "md_path": str(output_dir / "extracted.md")}


def extract_from_docx(input_path: Path, output_dir: Path) -> dict:
    """Extract text from DOCX file."""
    scripts_dir = Path(__file__).resolve().parent
    converter = scripts_dir / "source_to_md" / "doc_to_md.py"

    if converter.exists():
        import subprocess
        result = subprocess.run(
            [sys.executable, str(converter), str(input_path)],
            capture_output=True, text=True, cwd=str(scripts_dir.parent.parent)
        )
        if result.returncode == 0:
            md_path = input_path.with_suffix(".md")
            if md_path.exists():
                md_content = md_path.read_text(encoding="utf-8")
            else:
                md_content = result.stdout
        else:
            md_content = f"Error converting: {result.stderr[:500]}"
    else:
        try:
            from docx import Document
            doc = Document(str(input_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            md_content = "\n\n".join(paragraphs)
        except Exception as e:
            md_content = f"Error: {e}"

    output_dir.mkdir(parents=True, exist_ok=True)
    (output_dir / "extracted.md").write_text(md_content, encoding="utf-8")

    slides_data = _parse_markdown_to_slides(md_content)
    (output_dir / "slide_metadata.json").write_text(
        json.dumps(slides_data, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    return {"slides": slides_data, "total_slides": len(slides_data), "md_path": str(output_dir / "extracted.md")}


def extract_from_pdf(input_path: Path, output_dir: Path) -> dict:
    """Extract text + images from PDF file."""
    try:
        import fitz
        doc = fitz.open(str(input_path))
    except ImportError:
        return extract_from_docx(input_path, output_dir)

    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    slides_data = []
    all_text = []

    for page_idx, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        slide_images = []

        for img_idx, img in enumerate(page.get_images(full=True), 1):
            xref = img[0]
            base_image = doc.extract_image(xref)
            ext = base_image["ext"]
            img_name = f"slide_{page_idx:02d}_image_{img_idx:02d}.{ext}"
            img_path = images_dir / img_name
            with open(img_path, "wb") as f:
                f.write(base_image["image"])
            slide_images.append(img_name)

        title = text.split("\n")[0][:80] if text else f"Page {page_idx}"

        slides_data.append({
            "slide_number": page_idx,
            "title": title,
            "content_raw": text,
            "bullet_points": [],
            "table_data": [],
            "images_extracted": slide_images,
            "layout_hint": "content",
            "total_shapes": 0,
            "text_length": len(text),
        })

        all_text.append(f"## Page {page_idx}: {title}\n\n{text}")

    doc.close()

    output_dir.mkdir(parents=True, exist_ok=True)
    md_content = "\n\n---\n\n".join(all_text)
    (output_dir / "extracted.md").write_text(md_content, encoding="utf-8")
    (output_dir / "slide_metadata.json").write_text(
        json.dumps(slides_data, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    return {"slides": slides_data, "total_slides": len(slides_data), "md_path": str(output_dir / "extracted.md")}


def _parse_markdown_to_slides(md_content: str) -> list:
    """Parse markdown into slide objects by headers."""
    slides = []
    current_slide = None
    slide_num = 0

    for line in md_content.split("\n"):
        if re.match(r"^#{1,3}\s+", line):
            if current_slide and current_slide["content_raw"].strip():
                slides.append(current_slide)
            slide_num += 1
            title = re.sub(r"^#{1,3}\s+", "", line).strip()
            current_slide = {
                "slide_number": slide_num,
                "title": title,
                "content_raw": "",
                "bullet_points": [],
                "table_data": [],
                "images_extracted": [],
                "layout_hint": "content",
                "total_shapes": 0,
                "text_length": 0,
            }
        elif current_slide:
            current_slide["content_raw"] += line + "\n"
        else:
            slide_num += 1
            current_slide = {
                "slide_number": slide_num,
                "title": "Introduction",
                "content_raw": line + "\n",
                "bullet_points": [],
                "table_data": [],
                "images_extracted": [],
                "layout_hint": "title_slide",
                "total_shapes": 0,
                "text_length": 0,
            }

    if current_slide and current_slide["content_raw"].strip():
        current_slide["text_length"] = len(current_slide["content_raw"])
        slides.append(current_slide)

    if not slides:
        slides = [{
            "slide_number": 1, "title": "Content", "content_raw": md_content,
            "bullet_points": [], "table_data": [], "images_extracted": [],
            "layout_hint": "content", "total_shapes": 0, "text_length": len(md_content),
        }]

    return slides


def main():
    import argparse
    parser = argparse.ArgumentParser(description="Extract slide assets from documents")
    parser.add_argument("input_file", help="Input file (PPTX/DOCX/PDF)")
    parser.add_argument("-o", "--output", default=None, help="Output directory")
    args = parser.parse_args()

    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"Error: File not found: {input_path}", file=sys.stderr)
        return 1

    output_dir = Path(args.output) if args.output else Path("sources")
    output_dir.mkdir(parents=True, exist_ok=True)

    suffix = input_path.suffix.lower()
    if suffix == ".pptx":
        result = extract_from_pptx(input_path, output_dir)
    elif suffix == ".pdf":
        result = extract_from_pdf(input_path, output_dir)
    elif suffix in (".docx", ".doc"):
        result = extract_from_docx(input_path, output_dir)
    else:
        print(f"Unsupported format: {suffix}", file=sys.stderr)
        return 1

    print(f"[OK] Extracted {result['total_slides']} slides")
    print(f"[OK] Markdown: {result['md_path']}")
    print(f"[OK] Metadata: {output_dir / 'slide_metadata.json'}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
